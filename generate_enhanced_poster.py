#!/usr/bin/env python3
"""Generate an enhanced poster using template and existing content"""

import os
import sys
import json
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def load_existing_content():
    """Load the already generated content"""
    content_file = 'contents/<4o_4o>_fnode_bullet_point_content_0.json'
    with open(content_file, 'r') as f:
        return json.load(f)

def load_tree_split():
    """Load the tree split layout"""
    tree_file = 'tree_splits/<4o_4o>_fnode_tree_split_0.json'
    with open(tree_file, 'r') as f:
        return json.load(f)

def create_enhanced_poster():
    """Create enhanced poster with template"""
    
    print("Loading template...")
    prs = Presentation('template.pptx')
    slide = prs.slides[0]
    
    # Clear existing text
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()
    
    print("Loading content and layout...")
    content = load_existing_content()
    tree_split = load_tree_split()
    
    print(f"Creating poster with {len(content)} sections")
    
    # Add title
    title_content = content[0]
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(41), Inches(3)
    )
    tf = title_box.text_frame
    tf.clear()
    
    # Add title text
    p = tf.add_paragraph()
    for run_data in title_content['title'][0]['runs']:
        run = p.add_run()
        run.text = run_data['text']
        run.font.size = Pt(60)
        run.font.bold = run_data.get('bold', False)
    p.alignment = PP_ALIGN.CENTER
    
    # Add authors
    p = tf.add_paragraph()
    for run_data in title_content['textbox1'][0]['runs']:
        run = p.add_run()
        run.text = run_data['text']
        run.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER
    
    # Add affiliation
    p = tf.add_paragraph()
    for run_data in title_content['textbox1'][1]['runs']:
        run = p.add_run()
        run.text = run_data['text']
        run.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER
    
    # Create content panels with enhanced layout
    panel_configs = [
        {"section": 1, "x": 1, "y": 4.5, "w": 20, "h": 14},
        {"section": 2, "x": 22, "y": 4.5, "w": 20, "h": 14},
        {"section": 3, "x": 1, "y": 19, "w": 20, "h": 15},
        {"section": 4, "x": 22, "y": 19, "w": 20, "h": 7},
        {"section": 5, "x": 22, "y": 27, "w": 20, "h": 7},
    ]
    
    for config in panel_configs:
        if config["section"] >= len(content):
            continue
            
        section = content[config["section"]]
        
        # Add panel background
        panel = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(config["x"]), Inches(config["y"]),
            Inches(config["w"]), Inches(config["h"])
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(250, 250, 250)
        panel.line.color.rgb = RGBColor(47, 85, 151)
        panel.line.width = Pt(3)
        
        # Add section title
        title_box = slide.shapes.add_textbox(
            Inches(config["x"] + 0.5), Inches(config["y"] + 0.3),
            Inches(config["w"] - 1), Inches(2)
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.add_paragraph()
        
        for run_data in section['title'][0]['runs']:
            run = p.add_run()
            run.text = run_data['text']
            run.font.size = Pt(48)
            run.font.bold = True
        
        # Add bullet points
        content_y_offset = 2.5
        for textbox_key in ['textbox1', 'textbox2']:
            if textbox_key not in section:
                continue
                
            content_box = slide.shapes.add_textbox(
                Inches(config["x"] + 0.5), 
                Inches(config["y"] + content_y_offset),
                Inches(config["w"] - 1), 
                Inches((config["h"] - 3) / 2)
            )
            tf = content_box.text_frame
            tf.clear()
            
            for para_data in section[textbox_key]:
                p = tf.add_paragraph()
                p.level = para_data.get('level', 0)
                
                # Add bullet if specified
                if para_data.get('bullet', False):
                    p.text = "• "
                
                for run_data in para_data['runs']:
                    run = p.add_run()
                    run.text = run_data['text']
                    run.font.size = Pt(para_data.get('font_size', 24))
                    run.font.bold = run_data.get('bold', False)
            
            content_y_offset += (config["h"] - 3) / 2 + 0.5
    
    # Add sample figures
    print("Adding visual elements...")
    
    # Logo/decoration from template (preserve existing)
    shape_count = len(slide.shapes)
    
    # Save enhanced poster
    output_path = 'enhanced_template_poster.pptx'
    prs.save(output_path)
    print(f"✅ Enhanced poster saved to: {output_path}")
    print(f"   Total shapes: {len(slide.shapes)}")
    
    return output_path

if __name__ == "__main__":
    poster_path = create_enhanced_poster()
    print(f"\nPoster successfully generated!")
    print(f"Please open {poster_path} to view the result.")