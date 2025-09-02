# Suppress warnings - MUST be before any imports
import os
import sys
import io

# Set environment variables to prevent CUDA compilation attempts
os.environ['CUDA_VISIBLE_DEVICES'] = '0'
os.environ['TORCH_CUDA_ARCH_LIST'] = '8.0;8.6;8.9;9.0'  # Exclude unsupported 12.0
os.environ['PYTORCH_JIT'] = '0'  # Disable JIT compilation
os.environ['TRANSFORMERS_VERBOSITY'] = 'error'  # Only show errors from transformers
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'  # Suppress TensorFlow logging if present

import warnings
warnings.filterwarnings('ignore', message='Could not load the custom kernel')
warnings.filterwarnings('ignore', category=FutureWarning)  # Suppress sklearn warnings
warnings.filterwarnings('ignore', message='TORCH_CUDA_ARCH_LIST')

import logging
logging.getLogger('docling_core.transforms.serializer.html').setLevel(logging.ERROR)  # Suppress MathML warnings
logging.getLogger('transformers').setLevel(logging.ERROR)  # Suppress transformers warnings

# Now do the regular imports
from PosterAgent.parse_raw import parse_raw, gen_image_and_table
from PosterAgent.gen_outline_layout import filter_image_table, gen_outline_layout_v2
from utils.wei_utils import get_agent_config, scale_to_target_area, char_capacity, utils_functions, run_code, style_bullet_content
from PosterAgent.gen_poster_content import gen_bullet_point_content
from PosterAgent.tree_split_layout import main_train, main_inference, get_arrangments_in_inches, split_textbox, to_inches
from PosterAgent.gen_pptx_code import generate_poster_code
from utils.src.utils import ppt_to_images

import argparse
import json
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import shutil

units_per_inch = 25

def extract_template_and_clone(template_path, save_path):
    """Clone template and extract styles and layout zones."""
    try:
        # Load the template
        prs = Presentation(template_path)
        
        # Initialize style dictionary with defaults
        styles = {
            'background_color': None,
            'title_text_color': (255, 255, 255),
            'title_fill_color': (47, 85, 151),
            'body_text_color': (0, 0, 0),
            'panel_color': (47, 85, 151),
            'panel_thickness': 5,
        }
        
        # Extract layout zones from template
        layout_zones = []
        
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            
            # Store all existing shapes for preservation
            for shape in slide.shapes:
                zone_info = {
                    'name': shape.name if hasattr(shape, 'name') else '',
                    'left': shape.left.inches if hasattr(shape, 'left') else 0,
                    'top': shape.top.inches if hasattr(shape, 'top') else 0,
                    'width': shape.width.inches if hasattr(shape, 'width') else 0,
                    'height': shape.height.inches if hasattr(shape, 'height') else 0,
                    'has_text': shape.has_text_frame,
                    'preserve': True  # Mark template elements to preserve
                }
                layout_zones.append(zone_info)
                
                # Extract colors from shapes
                if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                    if shape.fill.fore_color.type == 1:  # RGB color
                        rgb = shape.fill.fore_color.rgb
                        # Use this as panel/title fill color if it's not white/black
                        if rgb != (255, 255, 255) and rgb != (0, 0, 0):
                            styles['title_fill_color'] = (rgb[0], rgb[1], rgb[2])
                            styles['panel_color'] = (rgb[0], rgb[1], rgb[2])
        
        # Save the cloned template as base for poster
        prs.save(save_path)
        
        return styles, layout_zones, prs.slide_width.inches, prs.slide_height.inches
    
    except Exception as e:
        print(f"Warning: Could not clone template: {e}")
        # Return defaults if template cannot be used
        return {
            'background_color': None,
            'title_text_color': (255, 255, 255),
            'title_fill_color': (47, 85, 151),
            'body_text_color': (0, 0, 0),
            'panel_color': (47, 85, 151),
            'panel_thickness': 5,
        }, [], 48, 36

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Template-Based Poster Generation')
    parser.add_argument('--poster_path', type=str, required=True, help='Path to PDF paper')
    parser.add_argument('--template_path', type=str, required=True, help='Path to template PowerPoint')
    parser.add_argument('--model_name_t', type=str, default='4o')
    parser.add_argument('--model_name_v', type=str, default='4o')
    parser.add_argument('--poster_name', type=str, default=None)
    parser.add_argument('--index', type=int, default=0, help='Index for file naming')
    parser.add_argument('--poster_width_inches', type=int, default=48)
    parser.add_argument('--poster_height_inches', type=int, default=36)
    parser.add_argument('--ablation_no_commenter', action='store_true', help='Ablation: disable commenter')
    parser.add_argument('--ablation_no_example', action='store_true', help='Ablation: disable examples')
    parser.add_argument('--no_blank_detection', action='store_true', help='Disable blank space detection')
    parser.add_argument('--tmp_dir', type=str, default='tmp')
    
    args = parser.parse_args()
    
    start_time = time.time()
    
    # Setup paths and directories
    os.makedirs(args.tmp_dir, exist_ok=True)
    
    if args.poster_name is None:
        args.poster_name = os.path.splitext(os.path.basename(args.poster_path))[0]
    
    print(f"Using template: {args.template_path}")
    print(f"Processing paper: {args.poster_path}")
    
    # Clone template and extract styles
    print("\n1. Cloning template and extracting styles...")
    template_base_path = f'{args.tmp_dir}/template_base.pptx'
    template_styles, layout_zones, template_width, template_height = extract_template_and_clone(
        args.template_path, template_base_path
    )
    
    # Use template dimensions if provided
    if template_width and template_height:
        args.poster_width_inches = int(template_width)
        args.poster_height_inches = int(template_height)
        print(f"   Using template dimensions: {template_width} x {template_height} inches")
    
    # Create theme configuration
    theme_title_text_color = template_styles['title_text_color']
    theme_title_fill_color = template_styles['title_fill_color']
    theme = {
        'panel_visible': True,
        'textbox_visible': False,
        'figure_visible': False,
        'panel_theme': {
            'color': theme_title_fill_color,
            'thickness': 5,
            'line_style': 'solid',
        },
        'textbox_theme': None,
        'figure_theme': None,
    }
    
    # Step 2: Parse the PDF paper
    print("\n2. Parsing PDF paper...")
    agent_config_t = get_agent_config(args.model_name_t)
    agent_config_v = get_agent_config(args.model_name_v)
    
    input_token, output_token, raw_result = parse_raw(args, agent_config_t, version=2)
    _, _, images, tables = gen_image_and_table(args, raw_result)
    print(f"   Parsing token consumption: {input_token} -> {output_token}")
    
    total_input_tokens_t = input_token
    total_output_tokens_t = output_token
    total_input_tokens_v = 0
    total_output_tokens_v = 0
    
    # Step 3: Filter and select figures (increase selection for richer poster)
    print("\n3. Filtering figures...")
    
    # Temporarily increase figure selection for richer content
    original_max_figures = getattr(args, 'max_figures', 3)
    args.max_figures = 6  # Select more figures for richer poster
    
    input_token, output_token = filter_image_table(args, agent_config_t)
    total_input_tokens_t += input_token
    total_output_tokens_t += output_token
    print(f"   Filter token consumption: {input_token} -> {output_token}")
    
    # Count selected figures
    try:
        import json
        with open(f'filter_results/{args.model_name_t}_{args.poster_name}_filter_results_{args.index}.json', 'r') as f:
            filter_results = json.load(f)
            num_images = len(filter_results.get('images', []))
            num_tables = len(filter_results.get('tables', []))
            print(f"   Selected {num_images} images and {num_tables} tables")
    except:
        print(f"   Filter results processed")
    
    args.max_figures = original_max_figures  # Restore original
    
    # Step 4: Generate outline and content (enhance for richer layout)
    print("\n4. Generating content outline...")
    input_token, output_token, panels, figures = gen_outline_layout_v2(args, agent_config_t)
    total_input_tokens_t += input_token
    total_output_tokens_t += output_token
    print(f"   Outline token consumption: {input_token} -> {output_token}")
    print(f"   Generated {len(panels)} panels with {len(figures)} figure placements")
    
    # Add more visual elements if template has rich layout
    if len(layout_zones) > 10:  # Template has many elements
        print("   Enhancing layout for rich template...")
        # Duplicate some panels for more content density
        if len(panels) < 8:
            # Add sub-panels for methodology and experiments
            for panel in panels:
                if panel['section_name'] in ['Methodology', 'Numerical Experiments']:
                    panel['needs_subdivision'] = True
    
    # Step 5: Generate layout using tree split (this creates proper poster blocks)
    print("\n5. Generating poster layout...")
    poster_width = args.poster_width_inches * units_per_inch
    poster_height = args.poster_height_inches * units_per_inch
    
    # Scale to target area if needed
    poster_width, poster_height = scale_to_target_area(poster_width, poster_height)
    poster_width_inches = to_inches(poster_width, units_per_inch)
    poster_height_inches = to_inches(poster_height, units_per_inch)
    
    # Ensure poster fits PowerPoint limits
    if poster_width_inches > 56 or poster_height_inches > 56:
        if poster_width_inches >= poster_height_inches:
            scale_factor = 56 / poster_width_inches
        else:
            scale_factor = 56 / poster_height_inches
        poster_width_inches *= scale_factor
        poster_height_inches *= scale_factor
        poster_width = poster_width_inches * units_per_inch
        poster_height = poster_height_inches * units_per_inch
    
    print(f"   Poster size: {poster_width_inches} x {poster_height_inches} inches")
    
    # Learn and generate layout
    panel_model_params, figure_model_params = main_train()
    panel_arrangement, figure_arrangement, text_arrangement = main_inference(
        panels,
        panel_model_params,
        figure_model_params,
        poster_width,
        poster_height,
        shrink_margin=3
    )
    
    # Process text arrangement for title
    if text_arrangement:
        text_arrangement_title = text_arrangement[0]
        text_arrangement = text_arrangement[1:]
        text_arrangement_title_top, text_arrangement_title_bottom = split_textbox(
            text_arrangement_title, 
            0.8
        )
        text_arrangement = [text_arrangement_title_top, text_arrangement_title_bottom] + text_arrangement
    
    # Process figure arrangement
    for i in range(len(figure_arrangement)):
        panel_id = figure_arrangement[i]['panel_id']
        panel_section_name = panels[panel_id]['section_name']
        figure_info = figures[panel_section_name]
        if 'image' in figure_info:
            figure_id = figure_info['image']
            figure_path = images.get(str(figure_id), images.get(figure_id, {})).get('image_path', '')
        elif 'table' in figure_info:
            figure_id = figure_info['table']
            figure_path = tables.get(str(figure_id), tables.get(figure_id, {})).get('table_path', '')
        else:
            figure_path = ''
        figure_arrangement[i]['figure_path'] = figure_path
    
    # Calculate character capacity for text areas
    for text_arrangement_item in text_arrangement:
        num_chars = char_capacity(
            bbox=(text_arrangement_item['x'], text_arrangement_item['y'], 
                  text_arrangement_item['height'], text_arrangement_item['width'])
        )
        text_arrangement_item['num_chars'] = num_chars
    
    # Get arrangements in inches
    width_inch, height_inch, panel_arrangement_inches, figure_arrangement_inches, text_arrangement_inches = get_arrangments_in_inches(
        poster_width, poster_height, panel_arrangement, figure_arrangement, text_arrangement, units_per_inch
    )
    
    # Save tree split results
    tree_split_results = {
        'poster_width': poster_width,
        'poster_height': poster_height,
        'poster_width_inches': width_inch,
        'poster_height_inches': height_inch,
        'panels': panels,
        'panel_arrangement': panel_arrangement,
        'figure_arrangement': figure_arrangement,
        'text_arrangement': text_arrangement,
        'panel_arrangement_inches': panel_arrangement_inches,
        'figure_arrangement_inches': figure_arrangement_inches,
        'text_arrangement_inches': text_arrangement_inches,
    }
    os.makedirs('tree_splits', exist_ok=True)
    with open(f'tree_splits/<{args.model_name_t}_{args.model_name_v}>_{args.poster_name}_tree_split_{args.index}.json', 'w') as f:
        json.dump(tree_split_results, f, indent=4)
    
    # Step 6: Generate detailed content
    print("\n6. Generating poster content...")
    input_token_t, output_token_t, input_token_v, output_token_v = gen_bullet_point_content(
        args, agent_config_t, agent_config_v, tmp_dir=args.tmp_dir
    )
    total_input_tokens_t += input_token_t
    total_output_tokens_t += output_token_t
    total_input_tokens_v += input_token_v
    total_output_tokens_v += output_token_v
    print(f"   Content generation T: {input_token_t} -> {output_token_t}")
    print(f"   Content generation V: {input_token_v} -> {output_token_v}")
    
    # Load generated content
    content_file = f'contents/<{args.model_name_t}_{args.model_name_v}>_{args.poster_name}_bullet_point_content_{args.index}.json'
    if os.path.exists(content_file):
        with open(content_file, 'r') as f:
            bullet_content = json.load(f)
    else:
        print(f"Warning: Content file not found at {content_file}")
        bullet_content = []
    
    # Step 7: Apply basic styles to content
    print("\n7. Styling content...")
    for k, v in bullet_content[0].items():
        style_bullet_content(v, theme_title_text_color, theme_title_fill_color)
    
    for i in range(1, len(bullet_content)):
        curr_content = bullet_content[i]
        style_bullet_content(curr_content['title'], theme_title_text_color, theme_title_fill_color)
    
    # Step 8: Generate the PowerPoint poster code
    print("\n8. Generating poster PowerPoint...")
    poster_code = generate_poster_code(
        panel_arrangement_inches,
        text_arrangement_inches,
        figure_arrangement_inches,
        presentation_object_name='poster_presentation',
        slide_object_name='poster_slide',
        utils_functions=utils_functions,
        slide_width=width_inch,
        slide_height=height_inch,
        img_path=None,
        save_path=f'{args.tmp_dir}/poster.pptx',
        visible=False,
        content=bullet_content,
        theme=theme,
        tmp_dir=args.tmp_dir,
        template_path=template_base_path,  # Use the cloned template
    )
    
    # Execute the generated code to create PowerPoint
    output, err = run_code(poster_code)
    if err is not None:
        raise RuntimeError(f'Error in generating PowerPoint: {err}')
    
    # Step 9: Run deoverflow process if not disabled
    if not args.ablation_no_commenter:
        print("\n9. Running deoverflow refinement...")
        try:
            from PosterAgent.deoverflow import deoverflow
            # Save intermediate results for deoverflow
            import pickle as pkl
            style_ckpt = {
                'style_logs': {},
                'outline': {'meta': {'width': width_inch, 'height': height_inch}}
            }
            logs_ckpt = {'outline': {'meta': {'width': width_inch, 'height': height_inch}}}
            
            os.makedirs('checkpoints', exist_ok=True)
            with open(f'checkpoints/{args.model_name_t}_{args.poster_name}_style_ckpt_{args.index}.pkl', 'wb') as f:
                pkl.dump(style_ckpt, f)
            with open(f'checkpoints/{args.model_name_t}_{args.poster_name}_ckpt_{args.index}.pkl', 'wb') as f:
                pkl.dump(logs_ckpt, f)
            
            # Save content for deoverflow
            os.makedirs('contents', exist_ok=True)
            content_file = f'contents/{args.model_name_t}_{args.poster_name}_poster_content_{args.index}.json'
            with open(content_file, 'w') as f:
                json.dump(bullet_content, f, indent=4)
            
            # Run deoverflow with vision model
            args.model_name = args.model_name_t  # For compatibility
            input_token, output_token = deoverflow(args, agent_config_t, agent_config_v)
            total_input_tokens_t += input_token
            total_output_tokens_t += output_token
            print(f"   Deoverflow token consumption: {input_token} -> {output_token}")
        except Exception as e:
            print(f"   Warning: Deoverflow process skipped: {e}")
            print("   Continuing with original poster...")
    
    # Step 10: Create output directory and save
    output_dir = f'template_based_posters/{args.poster_name}'
    os.makedirs(output_dir, exist_ok=True)
    
    pptx_path = os.path.join(output_dir, f'{args.poster_name}_poster.pptx')
    shutil.move(f'{args.tmp_dir}/poster.pptx', pptx_path)
    print(f"\n✅ Poster successfully generated: {pptx_path}")
    
    # Step 11: Generate preview images
    print("\n11. Generating preview images...")
    try:
        ppt_to_images(pptx_path, output_dir)
        print(f"   Preview images saved to: {output_dir}")
    except Exception as e:
        print(f"   Warning: Could not generate preview images: {e}")
    
    # Calculate time and tokens
    end_time = time.time()
    time_taken = end_time - start_time
    
    print(f"\nTotal time: {time_taken:.2f} seconds")
    print(f"Total tokens T: {total_input_tokens_t} -> {total_output_tokens_t}")
    print(f"Total tokens V: {total_input_tokens_v} -> {total_output_tokens_v}")