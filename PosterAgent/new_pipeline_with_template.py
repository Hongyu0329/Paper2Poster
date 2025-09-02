from PosterAgent.parse_raw import parse_raw, gen_image_and_table
from PosterAgent.gen_outline_layout import filter_image_table, gen_outline_layout_v2
from utils.wei_utils import get_agent_config, utils_functions, run_code, style_bullet_content, scale_to_target_area, char_capacity
from PosterAgent.tree_split_layout import main_train, main_inference, get_arrangments_in_inches, split_textbox, to_inches
from PosterAgent.gen_pptx_code import generate_poster_code
from utils.src.utils import ppt_to_images
from PosterAgent.gen_poster_content import gen_bullet_point_content
from utils.ablation_utils import no_tree_get_layout

import argparse
import json
import os
import time
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

units_per_inch = 25

def extract_template_styles(template_path):
    """Extract styles from a template PowerPoint file."""
    try:
        prs = Presentation(template_path)
        
        # Initialize style dictionary with defaults
        styles = {
            'background_color': None,
            'title_text_color': (255, 255, 255),
            'title_fill_color': (47, 85, 151),
            'body_text_color': (0, 0, 0),
            'panel_color': (47, 85, 151),
            'panel_thickness': 5,
            'font_name': 'Arial',
            'title_font_size': 44,
            'body_font_size': 20,
        }
        
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            
            # Try to extract background color
            if slide.background and hasattr(slide.background.fill, 'fore_color'):
                if slide.background.fill.fore_color.type == 1:  # RGB color
                    rgb = slide.background.fill.fore_color.rgb
                    styles['background_color'] = (rgb[0], rgb[1], rgb[2])
            
            # Extract text styles from shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            
                            # Get font properties
                            if run.font.name:
                                styles['font_name'] = run.font.name
                            
                            # Check if this looks like a title (larger font)
                            if run.font.size and run.font.size.pt > 30:
                                styles['title_font_size'] = run.font.size.pt
                                if run.font.color and run.font.color.type == 1:
                                    rgb = run.font.color.rgb
                                    styles['title_text_color'] = (rgb[0], rgb[1], rgb[2])
                            elif run.font.size:
                                styles['body_font_size'] = run.font.size.pt
                                if run.font.color and run.font.color.type == 1:
                                    rgb = run.font.color.rgb
                                    styles['body_text_color'] = (rgb[0], rgb[1], rgb[2])
                
                # Try to extract fill colors from shapes (potential panel colors)
                if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                    if shape.fill.fore_color.type == 1:  # RGB color
                        rgb = shape.fill.fore_color.rgb
                        # Use this as panel/title fill color if it's not white/black
                        if rgb != (255, 255, 255) and rgb != (0, 0, 0):
                            styles['title_fill_color'] = (rgb[0], rgb[1], rgb[2])
                            styles['panel_color'] = (rgb[0], rgb[1], rgb[2])
        
        print(f"Extracted template styles: {json.dumps(styles, indent=2)}")
        return styles
    
    except Exception as e:
        print(f"Warning: Could not extract all template styles: {e}")
        # Return default styles
        return {
            'background_color': None,
            'title_text_color': (255, 255, 255),
            'title_fill_color': (47, 85, 151),
            'body_text_color': (0, 0, 0),
            'panel_color': (47, 85, 151),
            'panel_thickness': 5,
            'font_name': 'Arial',
            'title_font_size': 44,
            'body_font_size': 20,
        }

def apply_template_to_poster(poster_path, template_styles):
    """Apply extracted template styles to the generated poster."""
    try:
        prs = Presentation(poster_path)
        
        # Apply to all slides
        for slide in prs.slides:
            # Apply background color if available
            if template_styles.get('background_color'):
                if slide.background.fill:
                    slide.background.fill.solid()
                    rgb = template_styles['background_color']
                    slide.background.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
            
            # Apply text styles to shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Apply font name
                            run.font.name = template_styles['font_name']
                            
                            # Determine if this is title or body text based on size
                            if run.font.size and run.font.size.pt > 30:
                                # Title text
                                rgb = template_styles['title_text_color']
                                run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                            else:
                                # Body text
                                rgb = template_styles['body_text_color']
                                run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                
                # Apply fill colors to shapes with solid fills
                if hasattr(shape, 'fill') and shape.fill.type == 1:
                    # Check if this might be a title/header shape
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        has_title = any(run.font.size and run.font.size.pt > 30 
                                       for para in shape.text_frame.paragraphs 
                                       for run in para.runs if run.font.size)
                        if has_title:
                            rgb = template_styles['title_fill_color']
                            shape.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        
        # Save the modified presentation
        prs.save(poster_path)
        print(f"Template styles applied to {poster_path}")
        
    except Exception as e:
        print(f"Warning: Could not apply all template styles: {e}")

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Poster Generation Pipeline with Template Support')
    parser.add_argument('--poster_path', type=str)
    parser.add_argument('--template_path', type=str, help='Path to template PowerPoint file')
    parser.add_argument('--model_name_t', type=str, default='4o')
    parser.add_argument('--model_name_v', type=str, default='4o')
    parser.add_argument('--index', type=int, default=0)
    parser.add_argument('--poster_name', type=str, default=None)
    parser.add_argument('--tmp_dir', type=str, default='tmp')
    parser.add_argument('--poster_width_inches', type=int, default=None)
    parser.add_argument('--poster_height_inches', type=int, default=None)
    parser.add_argument('--no_blank_detection', action='store_true', help='When overflow is severe, try this option.')
    parser.add_argument('--ablation_no_tree_layout', action='store_true', help='Ablation study: no tree layout')
    parser.add_argument('--ablation_no_commenter', action='store_true', help='Ablation study: no commenter')
    parser.add_argument('--ablation_no_example', action='store_true', help='Ablation study: no example')

    args = parser.parse_args()

    start_time = time.time()

    os.makedirs(args.tmp_dir, exist_ok=True)

    detail_log = {}

    # Extract template styles if template path is provided
    if args.template_path:
        template_styles = extract_template_styles(args.template_path)
        # Update theme colors based on template
        theme_title_text_color = template_styles['title_text_color']
        theme_title_fill_color = template_styles['title_fill_color']
    else:
        # Default theme
        theme_title_text_color = (255, 255, 255)
        theme_title_fill_color = (47, 85, 151)
        template_styles = None

    theme = {
        'panel_visible': True,
        'textbox_visible': False,
        'figure_visible': False,
        'panel_theme': {
            'color': theme_title_fill_color,
            'thickness': 5,
            'line_style': 'solid',
        },
        'textbox_theme': None,
        'figure_theme': None,
    }

    agent_config_t = get_agent_config(args.model_name_t)
    agent_config_v = get_agent_config(args.model_name_v)
    poster_name = args.poster_path.split('/')[-2].replace(' ', '_')
    if args.poster_name is None:
        args.poster_name = poster_name
    else:
        poster_name = args.poster_name
    meta_json_path = args.poster_path.replace('paper.pdf', 'meta.json')
    if args.poster_width_inches is not None and args.poster_height_inches is not None:
        poster_width = args.poster_width_inches * units_per_inch
        poster_height = args.poster_height_inches * units_per_inch
    elif os.path.exists(meta_json_path):
        meta_json = json.load(open(meta_json_path, 'r'))
        poster_width = meta_json['width']
        poster_height = meta_json['height']
    else:
        poster_width = 48 * units_per_inch
        poster_height = 36 * units_per_inch

    poster_width, poster_height = scale_to_target_area(poster_width, poster_height)
    poster_width_inches = to_inches(poster_width, units_per_inch)
    poster_height_inches = to_inches(poster_height, units_per_inch)

    if poster_width_inches > 56 or poster_height_inches > 56:
        # Work out which side is longer, then compute a single scale factor
        if poster_width_inches >= poster_height_inches:
            scale_factor = 56 / poster_width_inches
        else:
            scale_factor = 56 / poster_height_inches

        poster_width_inches  *= scale_factor
        poster_height_inches *= scale_factor

        # convert back to internal units
        poster_width  = poster_width_inches  * units_per_inch
        poster_height = poster_height_inches * units_per_inch

    print(f'Poster size: {poster_width_inches} x {poster_height_inches} inches')

    total_input_tokens_t, total_output_tokens_t = 0, 0
    total_input_tokens_v, total_output_tokens_v = 0, 0

    # Step 1: Parse the raw poster
    input_token, output_token, raw_result = parse_raw(args, agent_config_t, version=2)
    total_input_tokens_t += input_token
    total_output_tokens_t += output_token

    _, _, images, tables = gen_image_and_table(args, raw_result)

    print(f'Parsing token consumption: {input_token} -> {output_token}')

    detail_log['parser_in_t'] = input_token
    detail_log['parser_out_t'] = output_token


    # Step 2: Filter unnecessary images and tables
    input_token, output_token = filter_image_table(args, agent_config_t)
    total_input_tokens_t += input_token
    total_output_tokens_t += output_token
    print(f'Filter figures token consumption: {input_token} -> {output_token}')

    detail_log['filter_in_t'] = input_token
    detail_log['filter_out_t'] = output_token

    # Step 3: Generate outline
    input_token, output_token, panels, figures = gen_outline_layout_v2(args, agent_config_t)
    total_input_tokens_t += input_token
    total_output_tokens_t += output_token
    print(f'Outline token consumption: {input_token} -> {output_token}')

    detail_log['outline_in_t'] = input_token
    detail_log['outline_out_t'] = output_token

    if args.ablation_no_tree_layout:
        panel_arrangement, figure_arrangement, text_arrangement, input_token, output_token = no_tree_get_layout(
            poster_width, 
            poster_height, 
            panels, 
            figures, 
            agent_config_t
        )
        total_input_tokens_t += input_token
        total_output_tokens_t += output_token
        print(f'No tree layout token consumption: {input_token} -> {output_token}')
        detail_log['no_tree_layout_in_t'] = input_token
        detail_log['no_tree_layout_out_t'] = output_token
    else:

        # Step 4: Learn and generate layout
        panel_model_params, figure_model_params = main_train()

        panel_arrangement, figure_arrangement, text_arrangement = main_inference(
            panels,
            panel_model_params,
            figure_model_params,
            poster_width,
            poster_height,
            shrink_margin=3
        )

        text_arrangement_title = text_arrangement[0]
        text_arrangement = text_arrangement[1:]
        # Split the title textbox into two parts
        text_arrangement_title_top, text_arrangement_title_bottom = split_textbox(
            text_arrangement_title, 
            0.8
        )
        # Add the split textboxes back to the list
        text_arrangement = [text_arrangement_title_top, text_arrangement_title_bottom] + text_arrangement

    for i in range(len(figure_arrangement)):
        panel_id = figure_arrangement[i]['panel_id']
        panel_section_name = panels[panel_id]['section_name']
        figure_info = figures[panel_section_name]
        if 'image' in figure_info:
            figure_id = figure_info['image']
            if not figure_id in images:
                figure_path = images[str(figure_id)]['image_path']
            else:
                figure_path = images[figure_id]['image_path']
        elif 'table' in figure_info:
            figure_id = figure_info['table']
            if not figure_id in tables:
                figure_path = tables[str(figure_id)]['table_path']
            else:
                figure_path = tables[figure_id]['table_path']
        
        figure_arrangement[i]['figure_path'] = figure_path
        
    for text_arrangement_item in text_arrangement:
        num_chars = char_capacity(
            bbox=(text_arrangement_item['x'], text_arrangement_item['y'], text_arrangement_item['height'], text_arrangement_item['width'])
        )
        text_arrangement_item['num_chars'] = num_chars


    width_inch, height_inch, panel_arrangement_inches, figure_arrangement_inches, text_arrangement_inches = get_arrangments_in_inches(
        poster_width, poster_height, panel_arrangement, figure_arrangement, text_arrangement, 25
    )

    # Save to file
    tree_split_results = {
        'poster_width': poster_width,
        'poster_height': poster_height,
        'poster_width_inches': width_inch,
        'poster_height_inches': height_inch,
        'panels': panels,
        'panel_arrangement': panel_arrangement,
        'figure_arrangement': figure_arrangement,
        'text_arrangement': text_arrangement,
        'panel_arrangement_inches': panel_arrangement_inches,
        'figure_arrangement_inches': figure_arrangement_inches,
        'text_arrangement_inches': text_arrangement_inches,
    }
    os.makedirs('tree_splits', exist_ok=True)
    with open(f'tree_splits/<{args.model_name_t}_{args.model_name_v}>_{args.poster_name}_tree_split_{args.index}.json', 'w') as f:
        json.dump(tree_split_results, f, indent=4)

    # Step 5: Generate content
    input_token_t, output_token_t, input_token_v, output_token_v = gen_bullet_point_content(args, agent_config_t, agent_config_v, tmp_dir=args.tmp_dir)
    total_input_tokens_t += input_token
    total_output_tokens_t += output_token
    total_input_tokens_v += input_token_v
    total_output_tokens_v += output_token_v
    print(f'Content generation token consumption T: {input_token_t} -> {output_token_t}')
    print(f'Content generation token consumption V: {input_token_v} -> {output_token_v}')

    bullet_content = json.load(open(f'contents/<{args.model_name_t}_{args.model_name_v}>_{args.poster_name}_bullet_point_content_{args.index}.json', 'r'))

    detail_log['content_in_t'] = input_token_t
    detail_log['content_out_t'] = output_token_t
    detail_log['content_in_v'] = input_token_v
    detail_log['content_out_v'] = output_token_v
    
    # Step 6: Apply basic styles
    for k, v in bullet_content[0].items():
        style_bullet_content(v, theme_title_text_color, theme_title_fill_color)

    for i in range(1, len(bullet_content)):
        curr_content = bullet_content[i]
        style_bullet_content(curr_content['title'], theme_title_text_color, theme_title_fill_color)

    # Step 7: Generate the PowerPoint
    poster_code = generate_poster_code(
        panel_arrangement_inches,
        text_arrangement_inches,
        figure_arrangement_inches,
        presentation_object_name='poster_presentation',
        slide_object_name='poster_slide',
        utils_functions=utils_functions,
        slide_width=width_inch,
        slide_height=height_inch,
        img_path=None,
        save_path=f'{args.tmp_dir}/poster.pptx',
        visible=False,
        content=bullet_content,
        theme=theme,
        tmp_dir=args.tmp_dir,
    )

    output, err = run_code(poster_code)
    if err is not None:
        raise RuntimeError(f'Error in generating PowerPoint: {err}')
    
    # Step 7.5: Apply template styles if template was provided
    if args.template_path and template_styles:
        print("Applying template styles to generated poster...")
        apply_template_to_poster(f'{args.tmp_dir}/poster.pptx', template_styles)
    
    # Step 8: Create a folder in the output directory
    output_dir = f'<{args.model_name_t}_{args.model_name_v}>_generated_posters/{args.poster_path.replace("paper.pdf", "")}'
    os.makedirs(output_dir, exist_ok=True)

    # Step 9: Move poster.pptx to the output directory
    pptx_path = os.path.join(output_dir, f'{poster_name}.pptx')
    os.rename(f'{args.tmp_dir}/poster.pptx', pptx_path)
    print(f'Poster PowerPoint saved to {pptx_path}')
    # Step 10: Convert the PowerPoint to images
    ppt_to_images(pptx_path, output_dir)
    print(f'Poster images saved to {output_dir}')

    end_time = time.time()
    time_taken = end_time - start_time

    # log
    log_file = os.path.join(output_dir, 'log.json')
    with open(log_file, 'w') as f:
        log_data = {
            'input_tokens_t': total_input_tokens_t,
            'output_tokens_t': total_output_tokens_t,
            'input_tokens_v': total_input_tokens_v,
            'output_tokens_v': total_output_tokens_v,
            'time_taken': time_taken,
        }
        json.dump(log_data, f, indent=4)

    detail_log_file = os.path.join(output_dir, 'detail_log.json')
    with open(detail_log_file, 'w') as f:
        json.dump(detail_log, f, indent=4)