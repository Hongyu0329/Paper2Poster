#!/usr/bin/env python3
"""Test script for template-based poster generation"""

import os
import sys
sys.path.append('.')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def test_template_poster():
    """Create a simple poster using the template"""
    
    # Load template
    print("Loading template...")
    prs = Presentation('template.pptx')
    slide = prs.slides[0]
    
    print(f"Template has {len(slide.shapes)} shapes")
    
    # Clear existing text from template
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()
    
    # Add new content to specific positions
    print("Adding poster content...")
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(2), Inches(1),
        Inches(39), Inches(3)
    )
    title_box.text_frame.text = "FNODE: Flow-Matching for Multibody Systems"
    p = title_box.text_frame.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.alignment = 1  # Center
    
    # Add some content panels
    panels = [
        {"title": "Introduction", "x": 2, "y": 5, "w": 18, "h": 12},
        {"title": "Methodology", "x": 22, "y": 5, "w": 18, "h": 12},
        {"title": "Results", "x": 2, "y": 18, "w": 18, "h": 14},
        {"title": "Conclusions", "x": 22, "y": 18, "w": 18, "h": 14},
    ]
    
    for panel in panels:
        # Add panel background
        panel_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(panel["x"]), Inches(panel["y"]),
            Inches(panel["w"]), Inches(panel["h"])
        )
        panel_shape.fill.solid()
        panel_shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        panel_shape.line.color.rgb = RGBColor(47, 85, 151)
        panel_shape.line.width = Pt(2)
        
        # Add panel title
        title_box = slide.shapes.add_textbox(
            Inches(panel["x"] + 0.5), Inches(panel["y"] + 0.2),
            Inches(panel["w"] - 1), Inches(2)
        )
        tf = title_box.text_frame
        tf.clear()  # Clear default content
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = panel["title"]
        font = run.font
        font.size = Pt(48)
        font.bold = True
        # font.color.rgb = RGBColor(47, 85, 151)  # Skip color for now
        
        # Add sample content
        content_box = slide.shapes.add_textbox(
            Inches(panel["x"] + 0.5), Inches(panel["y"] + 2.5),
            Inches(panel["w"] - 1), Inches(panel["h"] - 3)
        )
        content_box.text_frame.text = f"• Sample content for {panel['title']}\n• More details here\n• Additional information"
        p = content_box.text_frame.paragraphs[0]
        p.font.size = Pt(24)
    
    # Save
    output_path = 'template_test_poster.pptx'
    prs.save(output_path)
    print(f"✅ Poster saved to: {output_path}")
    print(f"   Total shapes in poster: {len(slide.shapes)}")

if __name__ == "__main__":
    test_template_poster()